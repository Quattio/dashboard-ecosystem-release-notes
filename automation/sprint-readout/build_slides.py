"""Build the 2-slide Quatt-branded sprint readout deck (PowerPoint).

KNOWN-GOOD TEMPLATE — last verified against CW25 (16 Jun → 30 Jun 2026); the
slide-1 delivery-card layout was widened/heightened on 2026-06-29 after long
bodies overflowed the cards (see the delivery-card comment below — keep bodies
to ≤7 wrapped lines). The data constants below carry CW23 values as a worked
example. To reuse for a new sprint:
  1. Copy this file to a working copy, e.g. /tmp/build_cw<NN>_slides.py
  2. Replace the data constants with the new sprint's figures (steps 1-6):
       - TITLE_LINE / SUMMARY_LINE
       - TEAM_STATS (5 per-team cards — Embedded / Cloud / APP / Control / SW&I;
         set accent=True on the biggest contributor for the Neon pill)
       - DELIVERIES (3 cards; map to brand tints: Heating/All-E → Copper,
         Chill → Blue, Energy/HomeBattery → Green; bodies MUST carry per-team
         attribution)
       - RELEASES (fold non-tinted themes like CiC LED here)
       - TABLE_DATA (epic completion rows: (name, %done, sprint-Δ, is100))
       - SCOREBOARD + RISKS
  3. Run:  python3 /tmp/build_cw<NN>_slides.py        # → /tmp/sprint-cw<NN>-slides.pptx
  4. Slim: python3 assets/slim_pptx.py <in> ~/sprint-cw<NN>-slides.pptx
  5. The user drags ~/sprint-cw<NN>-slides.pptx into the 06 Reports Drive folder
     (do NOT MCP-upload the binary — it corrupts in transit).

Both slides are full dark mode (Forest Black #081412 background) — the deck
reads as a single specs-block presentation. Product distinction on slide 1
comes from the Copper-Dark / Blue-Dark / Green-Dark delivery-card backgrounds,
with Neon used selectively for accent (per-team stat pill, bullet markers,
percentages, 100%-Done table rows).

PER-TEAM REQUIREMENT (per SKILL.md): the five SW management teams
(Embedded Systems, Cloud / Backend, APP, Systems Control, SW&I) must each be
visible on slide 1 via the per-team stat strip, plus explicit attribution in
the delivery-card bodies and the footer. Cloud and APP come out of the SW&I
umbrella sprint but are first-class teams with their own products and releases;
never roll them up under "SW&I". The energyOS team is out of scope.

All tokens and rules from quatt-visual-branding/references/tokens.md apply.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# === Tokens (from quatt-visual-branding/references/tokens.md) ===
BLACK          = RGBColor(0x08, 0x14, 0x12)
DARKER         = RGBColor(0x28, 0x38, 0x36)   # cards on dark mode
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
NEON           = RGBColor(0xCC, 0xF8, 0x22)
GREY           = RGBColor(0xC8, 0xC8, 0xC8)
COPPER_DARK    = RGBColor(0x86, 0x59, 0x4E)   # Heating / All-Electric
BLUE_DARK      = RGBColor(0x39, 0x63, 0x6B)   # Chill (cooling)
GREEN_DARK     = RGBColor(0x57, 0x65, 0x54)   # Energy / HomeBattery
FONT = "Plus Jakarta Sans"

# === Per-sprint data constants (CW23 worked example — replace per sprint) ===
TITLE_LINE   = "Sprint CW23 — 2 Jun → 16 Jun 2026"
SUMMARY_LINE = ("5 teams  ·  72 tickets closed across 23 epics  ·  CiC + Cloud + App + "
                "Controller all shipped  ·  mid-sprint snapshot (day 6 of 14)")
TEAM_STATS = [
    ("7",  "Embedded Systems", False),
    ("34", "Cloud / Backend",  True),   # biggest contributor -> Neon pill
    ("20", "APP",              False),
    ("6",  "Systems Control",  False),
    ("5",  "SW&I (UX/QA)",     False),
]
DELIVERIES = [
    ("ALL-ELECTRIC", "All-E commissioning & Boost",
     "Systems Control cut the HC backup-heater pre/post-pump (~1.5 min/session) and swapped the 11-min Boost test for a ~30s 3-way-valve check; Cloud decoupled ALL-E diagnostics & admin-gated Boost; APP shipped the Shower-Minutes Boost UI.",
     COPPER_DARK),
    ("CHILL", "Chill commissioning & field stability",
     "Embedded fixed the RCP dongle TLS-PSK handshake; Systems Control fixed the startup reboot race; Cloud aggregated firmware-update state & cleared control bugs; APP fixed disconnected-Chill UI. Embedded → Controls → Cloud → APP.",
     BLUE_DARK),
    ("HOMEBATTERY", "HomeBattery insights",
     "Cloud shipped new earnings / metrics / PV-capacity endpoints + ingestion & processing; APP built the Earnings graph and insights polish. Full-stack Cloud → APP delivery; two insights epics at 93% and 100%.",
     GREEN_DARK),
]
# (lead, body) — fold non-tinted themes (e.g. CiC LED) into the right release.
RELEASES = [
    ("CiC 4.7.0", "(prep PR #762, 4 Jun) — Embedded: CiC LED UX rework, liveview persistence, LTE gating, RCP dongle TLS-PSK fix (+ Wireless Platform V2.11.X)"),
    ("Cloud v2.39.0", "(PR #3410, 3 Jun) — Cloud / Backend: Chill firmware-update state aggregation, ALL-E diagnostics decoupling"),
    ("App 1.58.0", "(PR #1873, 3 Jun) — APP: Boost Mode UI + Shower Minutes, HomeBattery insights screens"),
    ("Controller 6.9.0", "(active fixVersion) — Systems Control: Chill reboot fix, commissioning test changes"),
]
FOOTER = ("Source: Jira QPD sprints EMB / SW&I / Control - 26Q2 - CW23 — Embedded · Cloud / Backend · APP · Systems Control · SW&I"
          "  ·  GitHub: cic-yocto-builder, Quatt-cloud, quatt-mobile-app")

# (epic, %done, sprint Δ, is100) — sorted by % desc
TABLE_DATA = [
    ("App Re-architecture — App nav / Home / Product", "100%", "+3",  True),
    ("energyOS — HomeBattery historic activity/earnings", "100%", "+5", True),
    ("CHILL — App controls",                       "100%", "+5",  True),
    ("CHILL — Dongle SW enablement before launch",  "100%", "+1",  True),
    ("Automated E2E HomeBattery & Hybrid flows (QA)", "100%", "+1", True),
    ("CHILL — Preparation Pilot Phase 2",           "100%", "+1",  True),
    ("CHILL — Commissioning",                        "99%", "+2",  False),
    ("ODUv2 — Heat-pump OTA update",                 "97%", "+3",  False),
    ("CHILL — Multi-chill support",                  "97%", "+1",  False),
    ("energyOS — Expand & improve HomeBattery insights", "93%", "+10", False),
    ("CHILL — Field Test Stability Improvement",     "92%", "+9",  False),
    ("LED UX design for installers",                 "91%", "+2",  False),
    ("CHILL — Control Board OTA update",             "81%", "+1",  False),
    ("Migrate hybrid commissioning to new arch",     "80%", "+1",  False),
    ("CHILL — SW Improvements",                      "71%", "+2",  False),
    ("All-E — Boost Mode feature",                   "67%", "+2",  False),
    ("ALL-E — Commissioning & B2B enhancements",     "59%", "+4",  False),
    ("Installer Flow Improvement (Work Entity/CIC)", "36%", "+1",  False),
]
SCOREBOARD = [
    ("72", "tickets done across 5 teams & 23 epics"),
    ("6",  "epics fully completed (100% Done)"),
    ("14", "epics now at ≥80% — close-out candidates"),
    ("34", "biggest contributor — Cloud / Backend"),
]
RISKS = [
    ("Field:", "~9 customers reported blinking-red CiC LED on 4.6.0 (LTE down, LAN up). Fixed by the 4.7.0 LED UX rework; narrow corrective closed Won't-Do. Verify rollout closes the tickets."),
    ("Lagging:", "Installer Flow Improvement at 36% — newest & least-advanced active epic. Needs scoping."),
    ("In-flight:", "ALL-E Commissioning & B2B at 59% spans Controls + Cloud + APP; Boost→valve test swap has cross-team interface deps (ADR work)."),
    ("Snapshot:", "Captured day 6 of 14 — completion figures are a mid-sprint pulse, not final."),
]
BUCKETS_FOOTER = ("Ongoing buckets (no fixed scope, excluded from table): Production Incidents/Maintenance +10 (Cloud-heavy) · "
                  "Bug Fixing & Fleet Diagnostics · Tech Debt (Embedded + SW&I)")
SLIDE2_TITLE    = "Epic completion (% Done from Jira)  ·  Risks"
SLIDE2_SUBTITLE = "Sorted by completion. Δ = tickets closed against the epic this sprint. Mid-sprint — figures will rise before close."
OUT_PATH = "/tmp/sprint-cw23-slides.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_textbox(slide, x, y, w, h, anchor="top"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tf


def set_para(p, text, *, size=14, bold=False, color=BLACK, align=None, space_after=4, italic=False):
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT
    if align is not None:
        p.alignment = align
    p.space_after = Pt(space_after)
    return run


def add_card(slide, x, y, w, h, fill, border=None):
    box = slide.shapes.add_shape(5, x, y, w, h)  # ROUNDED_RECTANGLE = card-small
    if border is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = border
        box.line.width = Pt(1)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.shadow.inherit = False
    return box


def add_bullet(tf, text, *, size=13, bold_lead=None, color=BLACK, marker_color=None):
    p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
    marker = marker_color if marker_color is not None else color
    r1 = p.add_run(); r1.text = "• "
    r1.font.size = Pt(size); r1.font.color.rgb = marker; r1.font.bold = True; r1.font.name = FONT
    if bold_lead:
        r2 = p.add_run(); r2.text = bold_lead
        r2.font.size = Pt(size); r2.font.bold = True; r2.font.color.rgb = color; r2.font.name = FONT
        r3 = p.add_run(); r3.text = " " + text
        r3.font.size = Pt(size); r3.font.color.rgb = color; r3.font.name = FONT
    else:
        r2 = p.add_run(); r2.text = text
        r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = FONT
    p.space_after = Pt(4)


# ============= SLIDE 1 — DARK MODE (executive overview) =============
s1 = prs.slides.add_slide(BLANK)
bgall = s1.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
bgall.line.fill.background(); bgall.fill.solid(); bgall.fill.fore_color.rgb = BLACK
bgall.shadow.inherit = False

tf = add_textbox(s1, Inches(0.5), Inches(0.3), Inches(13.0), Inches(0.3))
set_para(tf.paragraphs[0], "QUATT · ENGINEERING READOUT", size=10, bold=True, color=NEON, space_after=0)
tf2 = add_textbox(s1, Inches(0.5), Inches(0.55), Inches(13.0), Inches(0.5))
set_para(tf2.paragraphs[0], TITLE_LINE, size=20, bold=True, color=WHITE, space_after=0)
tf3 = add_textbox(s1, Inches(0.5), Inches(0.95), Inches(13.0), Inches(0.25))
set_para(tf3.paragraphs[0], SUMMARY_LINE, size=11, color=GREY, space_after=0)

stat_y = Inches(1.5)
card_w = Inches(2.30); gap = Inches(0.20)
for i, (num, label, accent) in enumerate(TEAM_STATS):
    x = Inches(0.5) + i * (card_w + gap)
    add_card(s1, x, stat_y, card_w, Inches(1.0), DARKER)
    if accent:
        pill = s1.shapes.add_shape(5, x + Inches(0.20), stat_y + Inches(0.18), Inches(0.45), Inches(0.16))
        pill.line.fill.background(); pill.fill.solid(); pill.fill.fore_color.rgb = NEON
        pill.shadow.inherit = False
    tf = add_textbox(s1, x, stat_y, card_w, Inches(1.0), anchor="middle")
    set_para(tf.paragraphs[0], num, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=0)
    p = tf.add_paragraph()
    set_para(p, label, size=10, color=GREY, align=PP_ALIGN.CENTER, space_after=0)

tf = add_textbox(s1, Inches(0.5), Inches(2.62), Inches(12.3), Inches(0.4))
set_para(tf.paragraphs[0], "Three coordinated deliveries this sprint", size=15, bold=True, color=WHITE)

# Delivery cards: the rounded card is the visual container; the body textbox must
# fit INSIDE it. python-pptx does NOT auto-shrink overflowing text, so the body box
# is sized generously (1.32" tall) and the body font kept at 9.5pt — together ~7.7
# lines of headroom. Keep each DELIVERIES body to ≤7 wrapped lines (~370 chars at
# this width) or it will spill past the card edge. (Earlier 1.05"/10.5pt overflowed.)
for i, (eyebrow_t, title, body, fill) in enumerate(DELIVERIES):
    x = Inches(0.5 + i * 4.2)
    add_card(s1, x, Inches(3.02), Inches(4.05), Inches(2.34), fill)
    tf = add_textbox(s1, Inches(0.7 + i * 4.2), Inches(3.12), Inches(3.7), Inches(0.3))
    set_para(tf.paragraphs[0], eyebrow_t, size=9, bold=True, color=NEON, space_after=2)
    tf = add_textbox(s1, Inches(0.7 + i * 4.2), Inches(3.37), Inches(3.7), Inches(0.55))
    set_para(tf.paragraphs[0], title, size=14, bold=True, color=WHITE, space_after=2)
    tf = add_textbox(s1, Inches(0.7 + i * 4.2), Inches(3.96), Inches(3.7), Inches(1.32))
    set_para(tf.paragraphs[0], body, size=9.5, color=WHITE)

tf = add_textbox(s1, Inches(0.5), Inches(5.48), Inches(12.3), Inches(0.4))
set_para(tf.paragraphs[0], "Releases shipped / moved this sprint", size=15, bold=True, color=WHITE)

add_card(s1, Inches(0.5), Inches(5.84), Inches(12.3), Inches(1.28), DARKER)
rel_tf = add_textbox(s1, Inches(0.8), Inches(5.90), Inches(11.8), Inches(1.18))
for lead, body in RELEASES:
    add_bullet(rel_tf, body, bold_lead=lead, size=10.5, color=WHITE, marker_color=NEON)

fp = add_textbox(s1, Inches(0.5), Inches(7.2), Inches(12.3), Inches(0.25))
set_para(fp.paragraphs[0], FOOTER, size=8, color=GREY)


# ============= SLIDE 2 — DARK MODE (technical detail) =============
s2 = prs.slides.add_slide(BLANK)
bgall = s2.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
bgall.line.fill.background(); bgall.fill.solid(); bgall.fill.fore_color.rgb = BLACK
bgall.shadow.inherit = False

tf = add_textbox(s2, Inches(0.5), Inches(0.3), Inches(13.0), Inches(0.3))
set_para(tf.paragraphs[0], "QUATT · ENGINEERING READOUT", size=10, bold=True, color=NEON, space_after=0)
tf2 = add_textbox(s2, Inches(0.5), Inches(0.55), Inches(13.0), Inches(0.5))
set_para(tf2.paragraphs[0], SLIDE2_TITLE, size=20, bold=True, color=WHITE, space_after=0)
tf3 = add_textbox(s2, Inches(0.5), Inches(0.95), Inches(13.0), Inches(0.25))
set_para(tf3.paragraphs[0], SLIDE2_SUBTITLE, size=10, color=GREY, space_after=0)

table_x, table_y = Inches(0.5), Inches(1.35)
table_w = Inches(7.7)
rows = len(TABLE_DATA) + 1
tbl = s2.shapes.add_table(rows, 3, table_x, table_y, table_w, Inches(0.30 * rows)).table
tbl.columns[0].width = Inches(5.7)
tbl.columns[1].width = Inches(1.1)
tbl.columns[2].width = Inches(0.9)

for j, txt in enumerate(["Epic", "% Done", "Δ"]):
    cell = tbl.cell(0, j)
    cell.fill.solid(); cell.fill.fore_color.rgb = DARKER
    tf = cell.text_frame; tf.text = ""
    set_para(tf.paragraphs[0], txt, size=10, bold=True, color=NEON,
             align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT, space_after=0)
    cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)

for i, (name, pct, delta, done) in enumerate(TABLE_DATA, start=1):
    if done:
        bg_c = DARKER; text_color = NEON
    else:
        bg_c = BLACK if i % 2 == 1 else DARKER; text_color = WHITE
    for j, val in enumerate([name, pct, delta]):
        cell = tbl.cell(i, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = bg_c
        tf = cell.text_frame; tf.text = ""
        align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
        set_para(tf.paragraphs[0], val, size=9.5, bold=(j == 1 and done),
                 color=text_color, align=align, space_after=0)
        cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)

right_x = Inches(8.4); right_w = Inches(4.5)

add_card(s2, right_x, Inches(1.35), right_w, Inches(1.95), DARKER)
sb_tf = add_textbox(s2, right_x + Inches(0.25), Inches(1.5), right_w - Inches(0.5), Inches(1.8))
set_para(sb_tf.paragraphs[0], "Sprint scoreboard", size=13, bold=True, color=WHITE, space_after=6)
for lead, body in SCOREBOARD:
    add_bullet(sb_tf, body, bold_lead=lead, size=11.5, color=WHITE, marker_color=NEON)

add_card(s2, right_x, Inches(3.45), right_w, Inches(3.7), DARKER)
risk_tf = add_textbox(s2, right_x + Inches(0.25), Inches(3.6), right_w - Inches(0.5), Inches(3.45))
set_para(risk_tf.paragraphs[0], "Risks & call-outs", size=13, bold=True, color=WHITE, space_after=6)
for lead, body in RISKS:
    add_bullet(risk_tf, body, bold_lead=lead, size=10.5, color=WHITE, marker_color=GREY)

fp = add_textbox(s2, Inches(0.5), Inches(7.2), Inches(12.3), Inches(0.25))
set_para(fp.paragraphs[0], BUCKETS_FOOTER, size=8, color=GREY)

prs.save(OUT_PATH)
print("Saved", OUT_PATH, "slides:", len(prs.slides))
